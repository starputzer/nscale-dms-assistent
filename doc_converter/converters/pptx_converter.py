"""
PowerPoint-Konverter für die Dokumentenkonvertierungspipeline.
Konvertiert PowerPoint-Präsentationen (PPTX, PPT) in strukturierte Markdown-Dateien.
"""

import os
import re
from pathlib import Path
import logging
from typing import Dict, Any, List, Tuple, Optional, Set
from PIL import Image
import io
import pptx
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .base_converter import BaseConverter

class PowerPointConverter(BaseConverter):
    """Konvertiert PowerPoint-Präsentationen in strukturierte Markdown-Dateien"""
    
    def __init__(self, config: Optional[Dict[str, Any]] = None):
        """
        Initialisiert den PowerPoint-Konverter.
        
        Args:
            config: Konfigurationswörterbuch mit Konvertereinstellungen
        """
        super().__init__(config)
        
        # PowerPoint-spezifische Konfiguration
        self.extract_images = self.config.get('extract_images', True)
        self.extract_notes = self.config.get('extract_notes', True)
        self.extract_tables = self.config.get('extract_tables', True)
        self.max_image_size = self.config.get('max_image_size', 1024)
        self.preserve_animations = self.config.get('preserve_animations', False)
    
    def _convert_to_markdown(self, source_path: Path) -> Tuple[str, Dict[str, Any]]:
        """
        Konvertiert eine PowerPoint-Präsentation in Markdown.
        
        Args:
            source_path: Pfad zur PowerPoint-Datei
            
        Returns:
            Tuple aus (markdown_content, metadata)
        """
        self.logger.info(f"Starte Konvertierung von PowerPoint: {source_path}")
        
        try:
            # Öffne Präsentation
            presentation = Presentation(str(source_path))
            
            # Extrahiere Metadaten
            metadata = self._extract_pptx_metadata(presentation, source_path)
            
            # Erstelle Assets-Verzeichnis für Bilder
            assets_dir = source_path.parent / 'assets'
            assets_dir.mkdir(parents=True, exist_ok=True)
            
            # Konvertiere Folien
            markdown_parts = []
            markdown_parts.append(f"# {metadata['title']}\n")
            
            # Füge Metainformationen hinzu, falls vorhanden
            if 'author' in metadata and metadata['author']:
                markdown_parts.append(f"Autor: {metadata['author']}\n")
            
            # Inhaltsverzeichnis (wird später gefüllt)
            toc_placeholder = "[TOC_PLACEHOLDER]"
            markdown_parts.append(f"## Inhaltsverzeichnis\n\n{toc_placeholder}\n")
            
            # Konvertiere Folien
            slides_markdown, extracted_images = self._convert_slides(presentation, assets_dir)
            markdown_parts.append(slides_markdown)
            
            # Aktualisiere Metadaten mit Bildinformationen
            metadata['has_images'] = len(extracted_images) > 0
            metadata['extracted_images'] = extracted_images
            
            # Kombiniere alle Teile
            markdown_content = "\n".join(markdown_parts)
            
            # Generiere Inhaltsverzeichnis
            toc = self._generate_toc(markdown_content)
            markdown_content = markdown_content.replace(toc_placeholder, toc)
            
            self.logger.info(f"PowerPoint-Konvertierung abgeschlossen: {len(markdown_content)} Zeichen, {len(extracted_images)} Bilder")
            
            return markdown_content, metadata
            
        except Exception as e:
            self.logger.error(f"Fehler bei der PowerPoint-Konvertierung: {e}", exc_info=True)
            return f"# Konvertierung fehlgeschlagen\n\nFehler: {str(e)}", {
                'title': source_path.stem,
                'original_format': 'PPTX' if source_path.suffix.lower() == '.pptx' else 'PPT',
                'has_images': False,
                'error': str(e)
            }
    
    def _convert_slides(self, presentation: Presentation, assets_dir: Path) -> Tuple[str, List[Dict[str, Any]]]:
        """
        Konvertiert alle Folien einer Präsentation in Markdown.
        
        Args:
            presentation: PowerPoint-Präsentation
            assets_dir: Verzeichnis für extrahierte Bilder
            
        Returns:
            Tuple aus (markdown_content, extracted_images)
        """
        markdown_parts = []
        extracted_images = []
        
        for slide_idx, slide in enumerate(presentation.slides):
            # Erstelle Überschrift für Folie
            slide_num = slide_idx + 1
            
            # Versuche Titel der Folie zu extrahieren, falls vorhanden
            slide_title = self._extract_slide_title(slide)
            if slide_title:
                header = f"## Folie {slide_num}: {slide_title}"
            else:
                header = f"## Folie {slide_num}"
            
            markdown_parts.append(header)
            
            # Verarbeite Inhalte der Folie
            slide_content, slide_images = self._convert_slide_content(slide, slide_idx, assets_dir)
            markdown_parts.append(slide_content)
            
            # Füge extrahierte Bilder hinzu
            extracted_images.extend(slide_images)
            
            # Füge Notizen hinzu, falls vorhanden und aktiviert
            if self.extract_notes and slide.has_notes_slide and slide.notes_slide:
                notes_text = self._extract_notes_text(slide.notes_slide)
                if notes_text:
                    markdown_parts.append(f"\n### Notizen:\n\n{notes_text}\n")
        
        return "\n\n".join(markdown_parts), extracted_images
    
    def _extract_slide_title(self, slide) -> str:
        """
        Extrahiert den Titel einer Folie.
        
        Args:
            slide: PowerPoint-Folie
            
        Returns:
            Titel der Folie oder leerer String
        """
        for shape in slide.shapes:
            if hasattr(shape, "is_title") and shape.is_title:
                if shape.has_text_frame and shape.text_frame.text:
                    return shape.text_frame.text.strip()
        
        # Fallback: Suche nach großen Textboxen am oberen Rand
        potential_titles = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text:
                # Prüfe Position (oberes Drittel der Folie)
                if shape.top < slide.height / 3:
                    text = shape.text_frame.text.strip()
                    if text and len(text) < 100:  # Nicht zu lang für einen Titel
                        potential_titles.append((shape.top, text))
        
        # Sortiere nach Position (von oben nach unten)
        potential_titles.sort()
        
        # Verwende den obersten Text als Titel
        if potential_titles:
            return potential_titles[0][1]
        
        return ""
    
    def _convert_slide_content(self, slide, slide_idx: int, assets_dir: Path) -> Tuple[str, List[Dict[str, Any]]]:
        """
        Konvertiert den Inhalt einer Folie in Markdown.
        
        Args:
            slide: PowerPoint-Folie
            slide_idx: Index der Folie
            assets_dir: Verzeichnis für extrahierte Bilder
            
        Returns:
            Tuple aus (markdown_content, extracted_images)
        """
        content_parts = []
        extracted_images = []
        image_count = 0
        
        # Verarbeite alle Shapes auf der Folie
        for shape_idx, shape in enumerate(slide.shapes):
            shape_markdown = ""
            
            # Verarbeite Textframes
            if shape.has_text_frame:
                # Überspringe Shape, wenn es der Titel ist
                if hasattr(shape, "is_title") and shape.is_title:
                    continue
                
                text = self._process_text_frame(shape.text_frame)
                if text:
                    shape_markdown += text
            
            # Verarbeite Tabellen
            elif shape.has_table and self.extract_tables:
                table_markdown = self._process_table(shape.table)
                if table_markdown:
                    shape_markdown += f"\n{table_markdown}\n"
            
            # Verarbeite Bilder
            elif self.extract_images and (
                shape.shape_type == MSO_SHAPE_TYPE.PICTURE or 
                shape.shape_type == MSO_SHAPE_TYPE.LINKED_PICTURE or
                hasattr(shape, "image")
            ):
                try:
                    image_data = self._extract_image(shape, slide_idx, shape_idx, assets_dir)
                    if image_data:
                        image_count += 1
                        shape_markdown += f"\n![{image_data['alt']}]({image_data['path']})\n"
                        extracted_images.append(image_data)
                except Exception as e:
                    self.logger.warning(f"Fehler beim Extrahieren des Bildes auf Folie {slide_idx+1}, Shape {shape_idx}: {e}")
            
            # Gruppierte Shapes
            elif hasattr(shape, "shapes") and shape.shapes:
                try:
                    group_markdown = []
                    for group_shape in shape.shapes:
                        if group_shape.has_text_frame:
                            group_text = self._process_text_frame(group_shape.text_frame)
                            if group_text:
                                group_markdown.append(group_text)
                    
                    if group_markdown:
                        shape_markdown += "\n".join(group_markdown)
                except Exception as e:
                    self.logger.warning(f"Fehler bei der Verarbeitung einer Gruppe auf Folie {slide_idx+1}, Shape {shape_idx}: {e}")
            
            # Füge Shape-Inhalt zu Folie hinzu
            if shape_markdown:
                content_parts.append(shape_markdown)
        
        # Kombiniere alle Content-Teile
        if not content_parts:
            return "*Diese Folie enthält keinen extrahierbaren Text.*", extracted_images
        
        return "\n\n".join(content_parts), extracted_images
    
    def _process_text_frame(self, text_frame) -> str:
        """
        Verarbeitet einen TextFrame und konvertiert ihn in Markdown.
        
        Args:
            text_frame: PowerPoint-TextFrame
            
        Returns:
            Markdown-Text
        """
        # Wenn der TextFrame leer ist
        if not text_frame.text.strip():
            return ""
        
        paragraphs = []
        
        for paragraph in text_frame.paragraphs:
            # Ignoriere leere Absätze
            if not paragraph.text.strip():
                continue
            
            # Verarbeite Absatzformatierung
            text = paragraph.text.strip()
            
            # Erkennung von Listen
            if paragraph.level > 0:
                # Verschachtelte Liste
                indent = "  " * (paragraph.level - 1)
                paragraphs.append(f"{indent}* {text}")
            else:
                # Prüfe auf Aufzählungszeichen
                if text.startswith(("•", "◦", "▪", "▫", "◾", "◽")):
                    text = text[1:].strip()
                    paragraphs.append(f"* {text}")
                elif re.match(r"^\d+[\.\)]", text):
                    # Nummerierte Liste
                    text = re.sub(r"^\d+[\.\)]", "", text).strip()
                    paragraphs.append(f"1. {text}")
                else:
                    # Normaler Absatz mit Formatierung
                    formatted_text = self._apply_text_formatting(paragraph)
                    paragraphs.append(formatted_text)
        
        return "\n".join(paragraphs)
    
    def _apply_text_formatting(self, paragraph) -> str:
        """
        Wendet Textformatierungen auf einen Absatz an.
        
        Args:
            paragraph: PowerPoint-Paragraph
            
        Returns:
            Formatierter Markdown-Text
        """
        result = ""
        
        for run in paragraph.runs:
            text = run.text
            
            if not text:
                continue
            
            # Wende Formatierungen an
            if run.font.bold and run.font.italic:
                text = f"***{text}***"
            elif run.font.bold:
                text = f"**{text}**"
            elif run.font.italic:
                text = f"*{text}*"
            
            result += text
        
        return result
    
    def _process_table(self, table) -> str:
        """
        Konvertiert eine PowerPoint-Tabelle in Markdown.
        
        Args:
            table: PowerPoint-Tabelle
            
        Returns:
            Markdown-Tabelle
        """
        rows = []
        
        # Erstelle Header-Zeile
        header_row = []
        for cell in table.rows[0].cells:
            header_row.append(cell.text.strip() or " ")
        
        rows.append("| " + " | ".join(header_row) + " |")
        
        # Erstelle Trennzeile
        separator = "| " + " | ".join(["---"] * len(header_row)) + " |"
        rows.append(separator)
        
        # Erstelle Datenzeilen
        for row_idx, row in enumerate(table.rows):
            if row_idx == 0:  # Überspringe Header
                continue
            
            cells = []
            for cell in row.cells:
                cell_text = cell.text.strip().replace("\n", "<br>") or " "
                cells.append(cell_text)
            
            rows.append("| " + " | ".join(cells) + " |")
        
        return "\n".join(rows)
    
    def _extract_image(self, shape, slide_idx: int, shape_idx: int, assets_dir: Path) -> Optional[Dict[str, Any]]:
        """
        Extrahiert ein Bild aus einer PowerPoint-Shape.
        
        Args:
            shape: PowerPoint-Shape
            slide_idx: Index der Folie
            shape_idx: Index der Shape
            assets_dir: Verzeichnis für extrahierte Bilder
            
        Returns:
            Dictionary mit Bildinformationen oder None
        """
        if not hasattr(shape, "image"):
            return None
        
        # Extrahiere Bilddaten
        image_bytes = shape.image.blob
        if not image_bytes:
            return None
        
        # Generiere Dateinamen
        image_name = f"slide_{slide_idx+1}_image_{shape_idx+1}"
        
        try:
            # Öffne Bild mit PIL
            with Image.open(io.BytesIO(image_bytes)) as img:
                # Bestimme Dateiendung basierend auf Format
                file_ext = img.format.lower() if img.format else "png"
                
                # Skaliere große Bilder herunter
                width, height = img.size
                if width > self.max_image_size or height > self.max_image_size:
                    scale_factor = min(self.max_image_size / width, self.max_image_size / height)
                    new_width = int(width * scale_factor)
                    new_height = int(height * scale_factor)
                    img = img.resize((new_width, new_height), Image.LANCZOS)
                
                # Speichere Bild
                image_filename = f"{image_name}.{file_ext}"
                image_path = assets_dir / image_filename
                img.save(image_path)
                
                return {
                    'path': f"assets/{image_filename}",
                    'alt': f"Bild {slide_idx+1}-{shape_idx+1}",
                    'width': width,
                    'height': height,
                    'filename': image_filename,
                    'slide': slide_idx + 1,
                    'shape': shape_idx + 1
                }
        except Exception as e:
            self.logger.warning(f"Fehler beim Verarbeiten des Bildes: {e}")
            return None
    
    def _extract_notes_text(self, notes_slide) -> str:
        """
        Extrahiert den Text aus den Notizen einer Folie.
        
        Args:
            notes_slide: PowerPoint-NotesSlide
            
        Returns:
            Text der Notizen
        """
        notes_text = []
        
        for shape in notes_slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text and "Click to edit Master text styles" not in text:
                    notes_text.append(text)
        
        return "\n\n".join(notes_text)
    
    def _generate_toc(self, markdown_content: str) -> str:
        """
        Generiert ein Inhaltsverzeichnis basierend auf den Überschriften.
        
        Args:
            markdown_content: Markdown-Inhalt
            
        Returns:
            Inhaltsverzeichnis als Markdown
        """
        toc_entries = []
        
        # Finde alle Überschriften mit Regex
        slide_pattern = r'^## (Folie \d+:? ?.*?)$'
        
        for match in re.finditer(slide_pattern, markdown_content, re.MULTILINE):
            title = match.group(1)
            # Erstelle Anker-ID (GitHub-Style)
            anchor = title.lower().replace(" ", "-")
            anchor = re.sub(r'[^\w\-]', '', anchor)
            
            toc_entries.append(f"* [{title}](#{anchor})")
        
        if not toc_entries:
            return "*Kein Inhaltsverzeichnis verfügbar*"
        
        return "\n".join(toc_entries)
    
    def _extract_pptx_metadata(self, presentation: Presentation, source_path: Path) -> Dict[str, Any]:
        """
        Extrahiert Metadaten aus einer PowerPoint-Präsentation.
        
        Args:
            presentation: PowerPoint-Präsentation
            source_path: Pfad zur Präsentationsdatei
            
        Returns:
            Dictionary mit Metadaten
        """
        # Basis-Metadaten
        metadata = {
            'title': source_path.stem,
            'original_format': 'PPTX' if source_path.suffix.lower() == '.pptx' else 'PPT',
            'slide_count': len(presentation.slides),
            'has_images': False,
            'has_tables': False,
            'has_notes': False
        }
        
        # Core-Eigenschaften (falls verfügbar)
        if hasattr(presentation, 'core_properties'):
            core_props = presentation.core_properties
            if core_props.title:
                metadata['title'] = core_props.title
            
            if core_props.author:
                metadata['author'] = core_props.author
            
            if core_props.created:
                metadata['created'] = str(core_props.created)
            
            if core_props.modified:
                metadata['modified'] = str(core_props.modified)
            
            if core_props.subject:
                metadata['subject'] = core_props.subject
            
            if core_props.keywords:
                metadata['keywords'] = core_props.keywords
            
            if core_props.category:
                metadata['category'] = core_props.category
        
        # Präsentationseigenschaften
        if presentation.slide_width and presentation.slide_height:
            metadata['slide_width'] = presentation.slide_width
            metadata['slide_height'] = presentation.slide_height
        
        # Prüfe auf Notizen, Tabellen und Bilder
        notes_count = 0
        tables_count = 0
        images_count = 0
        
        for slide in presentation.slides:
            # Notizen
            if slide.has_notes_slide and slide.notes_slide:
                notes_count += 1
            
            # Bilder und Tabellen in Shapes
            for shape in slide.shapes:
                if shape.has_table:
                    tables_count += 1
                
                if (shape.shape_type == MSO_SHAPE_TYPE.PICTURE or 
                    shape.shape_type == MSO_SHAPE_TYPE.LINKED_PICTURE or
                    hasattr(shape, "image")):
                    images_count += 1
        
        metadata['has_notes'] = notes_count > 0
        metadata['has_tables'] = tables_count > 0
        metadata['has_images'] = images_count > 0
        metadata['notes_count'] = notes_count
        metadata['tables_count'] = tables_count
        metadata['images_count'] = images_count
        
        return metadata


if __name__ == "__main__":
    # Setup Logging
    logging.basicConfig(level=logging.INFO, 
                        format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
    
    # Test des PowerPoint-Konverters
    import sys
    
    if len(sys.argv) > 1:
        pptx_path = sys.argv[1]
        target_dir = sys.argv[2] if len(sys.argv) > 2 else "output"
        
        converter = PowerPointConverter()
        result = converter.convert(pptx_path, target_dir)
        
        if result['success']:
            print(f"Konvertierung erfolgreich: {result['target']}")
        else:
            print(f"Fehler bei der Konvertierung: {result.get('error', 'Unbekannter Fehler')}")
    else:
        print("Verwendung: python pptx_converter.py <pptx_datei> [ziel_verzeichnis]")