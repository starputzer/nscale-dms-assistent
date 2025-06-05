"""
Integrated Document Processor Module
Provides enhanced document processing capabilities for the RAG system
"""

import os
import logging
from typing import List, Dict, Any, Optional, Tuple
from datetime import datetime
import hashlib
import json
from pathlib import Path

logger = logging.getLogger(__name__)

class IntegratedDocumentProcessor:
    """Integrated processor for various document types"""
    
    def __init__(self):
        self.supported_formats = [
            '.txt', '.pdf', '.docx', '.doc', '.xlsx', '.xls', 
            '.pptx', '.ppt', '.md', '.json', '.csv', '.html'
        ]
        self.processing_stats = {
            'total_processed': 0,
            'successful': 0,
            'failed': 0,
            'by_type': {}
        }
        
    def process_document(self, file_path: str, metadata: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
        """Process a document and extract its content"""
        try:
            file_extension = Path(file_path).suffix.lower()
            
            if file_extension not in self.supported_formats:
                raise ValueError(f"Unsupported file format: {file_extension}")
                
            # Extract content based on file type
            content = self._extract_content(file_path, file_extension)
            
            # Generate document ID
            doc_id = self._generate_document_id(file_path, content)
            
            # Build document structure
            document = {
                'id': doc_id,
                'filename': os.path.basename(file_path),
                'path': file_path,
                'content': content,
                'file_type': file_extension,
                'size': os.path.getsize(file_path),
                'created_at': datetime.now().isoformat(),
                'metadata': metadata or {},
                'processing_info': {
                    'processor': 'integrated',
                    'version': '1.0',
                    'timestamp': datetime.now().isoformat()
                }
            }
            
            # Update stats
            self._update_stats(file_extension, success=True)
            
            return document
            
        except Exception as e:
            logger.error(f"Error processing document {file_path}: {str(e)}")
            self._update_stats(file_extension, success=False)
            raise
            
    def _extract_content(self, file_path: str, file_extension: str) -> str:
        """Extract content from file based on its type"""
        
        if file_extension in ['.txt', '.md']:
            return self._extract_text_content(file_path)
        elif file_extension == '.pdf':
            return self._extract_pdf_content(file_path)
        elif file_extension in ['.docx', '.doc']:
            return self._extract_word_content(file_path)
        elif file_extension in ['.xlsx', '.xls']:
            return self._extract_excel_content(file_path)
        elif file_extension in ['.pptx', '.ppt']:
            return self._extract_powerpoint_content(file_path)
        elif file_extension == '.json':
            return self._extract_json_content(file_path)
        elif file_extension == '.csv':
            return self._extract_csv_content(file_path)
        elif file_extension == '.html':
            return self._extract_html_content(file_path)
        else:
            raise ValueError(f"No extractor available for {file_extension}")
            
    def _extract_text_content(self, file_path: str) -> str:
        """Extract content from text files"""
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
            
    def _extract_pdf_content(self, file_path: str) -> str:
        """Extract content from PDF files"""
        try:
            import PyPDF2
            content = []
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    content.append(page.extract_text())
            return '\n'.join(content)
        except ImportError:
            logger.warning("PyPDF2 not installed, using fallback")
            return f"[PDF content from {file_path} - extraction not available]"
            
    def _extract_word_content(self, file_path: str) -> str:
        """Extract content from Word files"""
        try:
            from docx import Document
            doc = Document(file_path)
            content = []
            for paragraph in doc.paragraphs:
                content.append(paragraph.text)
            return '\n'.join(content)
        except ImportError:
            logger.warning("python-docx not installed, using fallback")
            return f"[Word content from {file_path} - extraction not available]"
            
    def _extract_excel_content(self, file_path: str) -> str:
        """Extract content from Excel files"""
        try:
            import openpyxl
            workbook = openpyxl.load_workbook(file_path)
            content = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                content.append(f"Sheet: {sheet_name}")
                for row in sheet.iter_rows(values_only=True):
                    row_text = '\t'.join(str(cell) if cell else '' for cell in row)
                    if row_text.strip():
                        content.append(row_text)
            return '\n'.join(content)
        except ImportError:
            logger.warning("openpyxl not installed, using fallback")
            return f"[Excel content from {file_path} - extraction not available]"
            
    def _extract_powerpoint_content(self, file_path: str) -> str:
        """Extract content from PowerPoint files"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            content = []
            for slide_num, slide in enumerate(prs.slides, 1):
                content.append(f"Slide {slide_num}:")
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text = shape.text.strip()
                        if text:
                            content.append(text)
            return '\n'.join(content)
        except ImportError:
            logger.warning("python-pptx not installed, using fallback")
            return f"[PowerPoint content from {file_path} - extraction not available]"
            
    def _extract_json_content(self, file_path: str) -> str:
        """Extract content from JSON files"""
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
            return json.dumps(data, indent=2, ensure_ascii=False)
            
    def _extract_csv_content(self, file_path: str) -> str:
        """Extract content from CSV files"""
        import csv
        content = []
        with open(file_path, 'r', encoding='utf-8') as f:
            csv_reader = csv.reader(f)
            for row in csv_reader:
                content.append('\t'.join(row))
        return '\n'.join(content)
        
    def _extract_html_content(self, file_path: str) -> str:
        """Extract content from HTML files"""
        try:
            from bs4 import BeautifulSoup
            with open(file_path, 'r', encoding='utf-8') as f:
                soup = BeautifulSoup(f.read(), 'html.parser')
                # Remove script and style elements
                for script in soup(["script", "style"]):
                    script.decompose()
                # Get text
                text = soup.get_text()
                # Break into lines and remove leading/trailing space
                lines = (line.strip() for line in text.splitlines())
                # Break multi-headlines into a line each
                chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
                # Drop blank lines
                text = '\n'.join(chunk for chunk in chunks if chunk)
                return text
        except ImportError:
            logger.warning("beautifulsoup4 not installed, using fallback")
            return f"[HTML content from {file_path} - extraction not available]"
            
    def _generate_document_id(self, file_path: str, content: str) -> str:
        """Generate a unique document ID"""
        hash_input = f"{file_path}:{content[:1000]}"
        return hashlib.sha256(hash_input.encode()).hexdigest()[:16]
        
    def _update_stats(self, file_extension: str, success: bool):
        """Update processing statistics"""
        self.processing_stats['total_processed'] += 1
        
        if success:
            self.processing_stats['successful'] += 1
        else:
            self.processing_stats['failed'] += 1
            
        if file_extension not in self.processing_stats['by_type']:
            self.processing_stats['by_type'][file_extension] = {
                'total': 0,
                'successful': 0,
                'failed': 0
            }
            
        self.processing_stats['by_type'][file_extension]['total'] += 1
        if success:
            self.processing_stats['by_type'][file_extension]['successful'] += 1
        else:
            self.processing_stats['by_type'][file_extension]['failed'] += 1
            
    def get_stats(self) -> Dict[str, Any]:
        """Get processing statistics"""
        return self.processing_stats.copy()
        
    def reset_stats(self):
        """Reset processing statistics"""
        self.processing_stats = {
            'total_processed': 0,
            'successful': 0,
            'failed': 0,
            'by_type': {}
        }
        
    def batch_process(self, file_paths: List[str], metadata: Optional[Dict[str, Any]] = None) -> List[Dict[str, Any]]:
        """Process multiple documents in batch"""
        results = []
        for file_path in file_paths:
            try:
                doc = self.process_document(file_path, metadata)
                results.append({
                    'status': 'success',
                    'document': doc
                })
            except Exception as e:
                results.append({
                    'status': 'error',
                    'file_path': file_path,
                    'error': str(e)
                })
        return results